import os
import fitz
import pandas as pd
from docx import Document
from pptx import Presentation

def extract_text(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == '.pdf':
            doc = fitz.open(filepath)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text.strip()

        elif ext in ['.docx', '.doc']:
            doc = Document(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
            return text.strip()

        elif ext in ['.xlsx', '.xls']:
            df = pd.read_excel(filepath)
            return df.to_string()

        elif ext == '.csv':
            df = pd.read_csv(filepath)
            return df.to_string()

        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()

        elif ext == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()

        return None
    except Exception as e:
        print(f"Error reading file: {e}")
        return None

def summarize_text(text, max_chars=4000):
    if len(text) > max_chars:
        return text[:max_chars] + "...[truncated]"
    return text